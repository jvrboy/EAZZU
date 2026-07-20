"""Presentations & visual storytelling tools.

Build slide decks as structured data and export to HTML reveal-style
markup, Markdown outline, JSON, or PPTX (when python-pptx is available).
"""
from __future__ import annotations

import html as _html
import json as _json
from pathlib import Path

_THEMES = {
    "corporate": {"bg": "#0f172a", "fg": "#f8fafc", "accent": "#38bdf8"},
    "warm": {"bg": "#fff7ed", "fg": "#7c2d12", "accent": "#f97316"},
    "forest": {"bg": "#052e16", "fg": "#ecfdf5", "accent": "#34d399"},
    "slate": {"bg": "#1e293b", "fg": "#e2e8f0", "accent": "#a78bfa"},
    "minimal": {"bg": "#ffffff", "fg": "#0f172a", "accent": "#0ea5e9"},
}

_TRANSITIONS = ["none", "fade", "slide", "push", "zoom", "flip"]
_ANIMATIONS = ["none", "fade-in", "slide-up", "slide-left", "zoom-in", "bounce"]


def _slide_html(slide, theme, index):
    title = _html.escape(slide.get("title", f"Slide {index + 1}"))
    bullets = slide.get("bullets", [])
    notes = slide.get("notes", "")
    body = [f"<h1>{title}</h1>"]
    if bullets:
        items = "".join(f"<li>{_html.escape(str(b))}</li>" for b in bullets)
        body.append(f"<ul>{items}</ul>")
    if slide.get("body"):
        body.append(f"<p>{_html.escape(slide['body'])}</p>")
    if notes:
        body.append(f"<aside class='notes'>{_html.escape(notes)}</aside>")
    return f"<section data-transition='{slide.get('transition','fade')}>{''.join(body)}</section>"


def _deck_html(deck):
    theme = _THEMES.get(deck.get("theme", "corporate"), _THEMES["corporate"])
    sections = "\n".join(_slide_html(s, theme, i) for i, s in enumerate(deck.get("slides", [])))
    return f"""<!doctype html><html><head><meta charset='utf-8'><title>{_html.escape(deck.get('title','Presentation'))}</title>
<style>body{{margin:0;background:{theme['bg']};color:{theme['fg']};font-family:system-ui}}
section{{display:none;padding:8vh 10vw;min-height:100vh}}section.active{{display:block}}
h1{{font-size:3rem;color:{theme['accent']}}}aside.notes{{display:none}}</style></head>
<body>{sections}<script>let i=0;const s=document.querySelectorAll('section');s[0].classList.add('active');
document.addEventListener('keydown',e=>{{if(e.key==='ArrowRight'||e.key===' '){{s[i].classList.remove('active');i=Math.min(i+1,s.length-1);s[i].classList.add('active')}}if(e.key==='ArrowLeft'){{s[i].classList.remove('active');i=Math.max(i-1,0);s[i].classList.add('active')}}}});</script></body></html>"""


def _deck_outline(deck):
    lines = [f"# {deck.get('title','Presentation')}", ""]
    for i, s in enumerate(deck.get("slides", []), 1):
        lines.append(f"{i}. {s.get('title', f'Slide {i}')}")
        for b in s.get("bullets", []):
            lines.append(f"   - {b}")
        if s.get("notes"):
            lines.append(f"   _notes: {s['notes']}")
    return "\n".join(lines)


def _export_deck(deck, fmt, path):
    if not path:
        return {"error": "path required"}
    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    if fmt == "html":
        p.write_text(_deck_html(deck), encoding="utf-8")
    elif fmt == "md":
        p.write_text(_deck_outline(deck), encoding="utf-8")
    elif fmt == "json":
        p.write_text(_json.dumps(deck, indent=2), encoding="utf-8")
    elif fmt == "pptx":
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation()
            for s in deck.get("slides", []):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = s.get("title", "")
                slide.placeholders[1].text = "\n".join(str(b) for b in s.get("bullets", []))
                if s.get("notes"):
                    slide.notes_slide.notes_text_frame.text = s["notes"]
            prs.save(str(p))
        except Exception as e:
            p.write_text(_deck_outline(deck), encoding="utf-8")
            return {"path": str(p), "fallback": "md", "error": f"python-pptx unavailable: {e}"}
    else:
        return {"error": f"unsupported format: {fmt}"}
    return {"path": str(p), "format": fmt, "bytes": p.stat().st_size}


def _design_suggest(slide):
    ideas = []
    n_bullets = len(slide.get("bullets", []))
    if n_bullets > 6:
        ideas.append("Split into two slides — more than 6 bullets reduces readability.")
    if len(slide.get("title", "")) > 60:
        ideas.append("Shorten the title to under 60 characters for impact.")
    if not slide.get("bullets") and not slide.get("body"):
        ideas.append("Add content or convert to a section divider.")
    if not ideas:
        ideas.append("Layout looks balanced. Consider adding an image or chart.")
    return ideas


def _outline_to_deck(outline, theme="corporate"):
    deck = {"title": "Presentation", "theme": theme, "slides": []}
    current = None
    for line in outline.splitlines():
        if line.startswith("# ") and not line.startswith("## "):
            deck["title"] = line[2:].strip()
        elif line.startswith("## "):
            if current:
                deck["slides"].append(current)
            current = {"title": line[3:].strip(), "bullets": [], "notes": "", "transition": "fade", "animation": "fade-in"}
        elif line.strip().startswith("- ") and current is not None:
            current["bullets"].append(line.strip()[2:])
        elif line.strip().startswith("_notes:") and current is not None:
            current["notes"] = line.strip()[7:].strip()
    if current:
        deck["slides"].append(current)
    return deck


TOOLS: list[dict] = [
    {"name": "slides_new_deck", "description": "Create a new slide deck structure with a title, theme, and empty slides list.", "params": {"title": "string", "theme": "string"}, "run": lambda a: {"deck": {"title": a.get("title", "Untitled"), "theme": a.get("theme", "corporate"), "slides": []}}},
    {"name": "slides_add_slide", "description": "Add a slide to a deck with title, bullets, body, notes, transition, and animation.", "params": {"deck": "object", "title": "string", "bullets": "list", "body": "string", "notes": "string", "transition": "string", "animation": "string"}, "run": lambda a: {"deck": {**a.get("deck", {}), "slides": a.get("deck", {}).get("slides", []) + [{"title": a.get("title", ""), "bullets": a.get("bullets", []), "body": a.get("body", ""), "notes": a.get("notes", ""), "transition": a.get("transition", "fade"), "animation": a.get("animation", "fade-in")}]}}},
    {"name": "slides_list_themes", "description": "List available presentation themes and their color palettes.", "params": {}, "run": lambda a: {"themes": _THEMES, "transitions": _TRANSITIONS, "animations": _ANIMATIONS}},
    {"name": "slides_design_suggest", "description": "Get automatic design suggestions for a slide based on its content.", "params": {"slide": "object"}, "run": lambda a: {"suggestions": _design_suggest(a.get("slide", {}))}},
    {"name": "slides_export", "description": "Export a deck to html (reveal-style), md (outline), json, or pptx (falls back to md if python-pptx unavailable).", "params": {"deck": "object", "format": "string", "path": "string"}, "run": lambda a: _export_deck(a.get("deck", {}), a.get("format", "html"), a.get("path", ""))},
    {"name": "slides_outline", "description": "Produce a Markdown outline of a deck (titles, bullets, notes).", "params": {"deck": "object"}, "run": lambda a: {"outline": _deck_outline(a.get("deck", {}))}},
    {"name": "slides_from_outline", "description": "Convert a Markdown outline (lines starting with # for title, ## for slide, - for bullets) into a deck structure.", "params": {"outline": "string", "theme": "string"}, "run": lambda a: {"deck": _outline_to_deck(a.get("outline", ""), a.get("theme", "corporate"))}},
]
